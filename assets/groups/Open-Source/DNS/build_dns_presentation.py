#!/usr/bin/env python3
"""
Cascade STEAM DNS Presentation Builder — v2 (enhanced visuals)
Theme: brand SVG colors | teal #34b0bf | orange #d46329 | navy #0a2c3f
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand Colors ──────────────────────────────────────────────────────────────
BG        = RGBColor(0x0a, 0x30, 0x42)
TEAL      = RGBColor(0x34, 0xb0, 0xbf)
ORANGE    = RGBColor(0xd4, 0x63, 0x29)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BOX_DARK  = RGBColor(0x06, 0x1d, 0x29)
BOX_MID   = RGBColor(0x0e, 0x3e, 0x54)
BOX_ORANGE= RGBColor(0x3a, 0x1a, 0x06)
GRAY      = RGBColor(0xcc, 0xdd, 0xe6)
TOPO_CLR  = RGBColor(0x0f, 0x3e, 0x54)   # Barely visible topo lines
GREEN     = RGBColor(0x27, 0xae, 0x60)
RED_CLR   = RGBColor(0xeb, 0x57, 0x57)

# ── Auto-Shape Type IDs (confirmed by enumeration) ────────────────────────────
SH_RECT      = 1
SH_ROUNDED   = 5
SH_OVAL      = 9
SH_R_ARROW   = 33
SH_PENTAGON  = 51
SH_CHEVRON   = 52

# ── Paths ─────────────────────────────────────────────────────────────────────
LOGO = os.path.expanduser(
    "~/Obsidian/CascadeSTEAM/assets/images/Cascade_STEAM_horizontal_logo_primary_darkBG.png"
)
OUT_DIR  = os.path.dirname(os.path.abspath(__file__))
PPTX_OUT = os.path.join(OUT_DIR, "DNS_Presentation_CascadeSTEAM.pptx")
MD_OUT   = os.path.join(OUT_DIR, "DNS_Presentation.md")

TOTAL = 18

# ── Topographic line presets (x, y, w) at h=0.022" ──────────────────────────
TOPO_SPECS = [
    (0.0,  0.65,  9.8), (1.8,  1.30,  7.5), (0.0,  1.95, 11.0),
    (0.5,  2.60,  6.2), (3.2,  2.72,  4.8), (0.0,  3.30, 13.3),
    (1.2,  3.92,  8.0), (0.0,  4.55, 10.5), (2.5,  4.70,  5.0),
    (0.0,  5.28, 13.3), (0.8,  5.90,  7.2), (4.5,  6.05,  4.0),
    (0.0,  6.65,  6.5), (5.0,  6.80,  4.8), (0.0,  7.28, 13.3),
]

# ── Core Helpers ──────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def shape(slide, stype, x, y, w, h, fill=None, line_color=None, line_w=Pt(1.5)):
    sp = slide.shapes.add_shape(stype, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    if line_color:
        sp.line.color.rgb = line_color
        sp.line.width = line_w
    else:
        sp.line.fill.background()
    return sp

def rect(slide, x, y, w, h, fill=None, line_color=None, line_w=Pt(1.5)):
    return shape(slide, SH_RECT, x, y, w, h, fill=fill, line_color=line_color, line_w=line_w)

def oval(slide, x, y, w, h, fill=None, line_color=None, line_w=Pt(1.5)):
    return shape(slide, SH_OVAL, x, y, w, h, fill=fill, line_color=line_color, line_w=line_w)

def rarrow(slide, x, y, w, h, fill=ORANGE):
    return shape(slide, SH_R_ARROW, x, y, w, h, fill=fill)

def txt(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tb

def code_box(slide, code_text, x, y, w, h):
    rect(slide, x, y, w, h, fill=BOX_DARK, line_color=TEAL, line_w=Pt(1))
    txt(slide, code_text, x + 0.12, y + 0.1, w - 0.24, h - 0.2,
        size=13, color=TEAL, font="Courier New")

def logo_img(slide, x, y, w=2.4):
    slide.shapes.add_picture(LOGO, Inches(x), Inches(y), width=Inches(w))

def header_logo(slide):
    logo_img(slide, 10.9, 0.2, w=2.2)

def footer(slide, num):
    txt(slide, "cascadesteam.org", 0.3, 7.1, 4, 0.35, size=11, color=TEAL)
    txt(slide, f"{num} / {TOTAL}", 11.8, 7.1, 1.5, 0.35, size=11,
        color=GRAY, align=PP_ALIGN.RIGHT)

def title_bar(slide, title, num, section_label=""):
    txt(slide, title, 0.4, 0.35, 10.3, 0.85, size=32, bold=True, color=WHITE)
    if section_label:
        txt(slide, section_label, 0.4, 1.05, 6, 0.35, size=14, color=ORANGE)
    rect(slide, 0.4, 1.15 if not section_label else 1.38, 12.5, 0.04, fill=ORANGE)
    header_logo(slide)
    footer(slide, num)

def topo_lines(slide):
    """Subtle topographic contour line decoration."""
    for (x, y, w) in TOPO_SPECS:
        rect(slide, x, y, w, 0.022, fill=TOPO_CLR)

# ── Globe Icon ────────────────────────────────────────────────────────────────
def draw_globe(slide, cx, cy, r):
    """Simple globe from concentric ovals and axis lines."""
    # Outer circle
    oval(slide, cx-r, cy-r, r*2, r*2, line_color=TEAL, line_w=Pt(1.5))
    # Latitude lines
    for frac, yscale in [(0.3, 0.5), (0.55, 0.85), (0.75, 0.62)]:
        oval(slide, cx-r, cy-r*yscale, r*2, r*yscale*2, line_color=TEAL, line_w=Pt(0.75))
    # Vertical axis (longitude)
    rect(slide, cx-0.012, cy-r, 0.024, r*2, fill=TEAL)
    rect(slide, cx-r, cy-0.012, r*2, 0.024, fill=TEAL)

# ── Lock Icon ─────────────────────────────────────────────────────────────────
def draw_lock(slide, x, y, w, h, color=TEAL):
    """Padlock shape from rectangle + oval shackle."""
    body_h = h * 0.58
    shackle_h = h * 0.50
    shackle_w = w * 0.5
    # Body
    shape(slide, SH_ROUNDED, x, y + h * 0.42, w, body_h, fill=color)
    # Shackle (U-shape approximated with oval outline)
    oval(slide, x + w*0.25, y, shackle_w, shackle_h, line_color=color, line_w=Pt(3))
    # Keyhole
    oval(slide, x + w*0.4, y + h*0.58, w*0.2, w*0.2,
         fill=BOX_DARK if color == TEAL else BG)

# ── DNS Hierarchy Tree ────────────────────────────────────────────────────────
def draw_dns_tree(slide, x, y):
    """Root → TLD → Zone hierarchy diagram."""
    node_w, node_h = 2.0, 0.52
    gap_x = 0.35
    col_w = node_w + gap_x
    # Level labels
    levels = [
        [(". (Root)", ORANGE, x + col_w, y)],
        [(".org TLD", TEAL,   x + col_w, y + 1.0)],
        [("cascadesteam.org", BOX_MID, x, y + 2.0),
         ("ns1.hover.com",   BOX_MID, x + col_w*1.2 + 0.4, y + 2.0)],
    ]
    nodes = []
    for level in levels:
        for label, color, nx, ny in level:
            rect(slide, nx, ny, node_w, node_h, fill=color,
                 line_color=TEAL, line_w=Pt(0.75))
            txt(slide, label, nx + 0.1, ny + 0.08, node_w - 0.2, node_h - 0.1,
                size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            nodes.append((nx, ny, node_w, node_h))
    # Connector lines (approximated with thin rects)
    # Root → TLD
    cx0 = x + col_w + node_w/2
    rect(slide, cx0 - 0.012, y + node_h, 0.024, 1.0 - node_h, fill=TEAL)
    # TLD → zone nodes
    cx1 = x + col_w + node_w/2
    cy1 = y + 1.0 + node_h
    rect(slide, cx1 - 0.012, cy1, 0.024, 0.4, fill=TEAL)
    # Horizontal branch
    lx = x + node_w/2
    rx = x + col_w*1.2 + 0.4 + node_w/2
    branch_y = cy1 + 0.4
    rect(slide, lx, branch_y - 0.012, rx - lx, 0.024, fill=TEAL)
    # Down to left node
    rect(slide, lx - 0.012, branch_y, 0.024, 2.0 - (branch_y - y), fill=TEAL)
    # Down to right node
    rect(slide, rx - 0.012, branch_y, 0.024, 2.0 - (branch_y - y), fill=TEAL)


# ── Slide 01 — Title ──────────────────────────────────────────────────────────
def slide_01_title(prs):
    s = blank_slide(prs)
    set_bg(s)
    topo_lines(s)

    # Left edge accent bar
    rect(s, 0, 0, 0.12, 7.5, fill=ORANGE)

    # Decorative globe in top-right (subtle)
    draw_globe(s, 10.8, 2.8, 2.0)

    # "DNS" massive title
    txt(s, "DNS", 0.5, 0.9, 8, 1.7, size=110, bold=True, color=WHITE)
    # Teal accent line
    rect(s, 0.5, 2.65, 8.5, 0.05, fill=TEAL)
    # Subtitle with orange arrow accent
    rarrow(s, 0.5, 2.78, 0.5, 0.42)
    txt(s, "More Than Just a Lookup Table",
        1.12, 2.75, 8.5, 0.58, size=28, bold=True, color=TEAL)
    txt(s, "Domain Name System — A deep dive into the internet's phone book,\n"
           "directory, security layer, and so much more.",
        0.5, 3.45, 8.5, 1.0, size=18, color=GRAY)
    txt(s, "Cascade STEAM — Technology Education Series",
        0.5, 4.6, 8.5, 0.45, size=14, color=ORANGE)

    logo_img(s, 10.1, 6.1, w=3.0)
    footer(s, 1)
    return s


# ── Section Slide Template ────────────────────────────────────────────────────
def _section_slide(prs, number, title, subtitle, num, icon_fn=None):
    s = blank_slide(prs)
    set_bg(s)
    topo_lines(s)

    rect(s, 0, 0, 0.5, 7.5, fill=ORANGE)
    rect(s, 0.5, 0, 0.08, 7.5, fill=TEAL)

    # Section number (watermark)
    txt(s, number, 0.9, 1.5, 4, 3.5, size=180, bold=True,
        color=RGBColor(0x0f, 0x3a, 0x50))
    # Overlay real number
    txt(s, number, 0.85, 1.8, 3, 1.5, size=80, bold=True, color=ORANGE)

    # Title and subtitle
    txt(s, title,    0.85, 3.35, 9,   1.0, size=48, bold=True, color=WHITE)
    txt(s, subtitle, 0.85, 4.45, 9,   0.75, size=20, color=TEAL)

    # Orange arrow accent
    rarrow(s, 0.85, 5.3, 1.5, 0.45)

    # Section-specific icon in right area
    if icon_fn:
        icon_fn(s)

    header_logo(s)
    footer(s, num)
    return s


def slide_03_section_what_is_dns(prs):
    def icon(s):
        draw_globe(s, 10.5, 3.2, 1.9)
    return _section_slide(prs, "01", "What is DNS?",
                          "The basics — names, IPs, and the phone book of the internet",
                          3, icon_fn=icon)

def slide_06_section_records(prs):
    def icon(s):
        # Stacked "record card" visual
        for i in range(3):
            off = i * 0.18
            rect(s, 9.0 + off, 2.0 + off, 3.5, 2.2, fill=BOX_MID,
                 line_color=TEAL, line_w=Pt(1))
            for li in range(4):
                rect(s, 9.2 + off, 2.5 + off + li*0.42, 3.0, 0.12, fill=TOPO_CLR)
    return _section_slide(prs, "02", "DNS Record Types",
                          "The Building Blocks — every record has a job",
                          6, icon_fn=icon)

def slide_13_section_wild(prs):
    def icon(s):
        # Network spiderweb: center circle + radiating spokes to service nodes
        cx, cy, spoke = 10.5, 3.5, 1.65
        # Outer nodes
        import math
        services = ["Web", "Mail", "CDN", "VoIP", "API"]
        for i, svc in enumerate(services):
            angle = math.radians(-20 + i * 52)
            nx = cx + spoke * math.cos(angle)
            ny = cy + spoke * math.sin(angle)
            oval(s, nx-0.38, ny-0.28, 0.76, 0.56, fill=BOX_MID, line_color=TEAL, line_w=Pt(1))
            txt(s, svc, nx-0.38, ny-0.18, 0.76, 0.36, size=10, bold=True,
                color=TEAL, align=PP_ALIGN.CENTER)
            # Spoke line
            lx = min(cx, nx) + abs(cx-nx)*0.22
            lw = abs(cx - nx) * 0.56
            ly = min(cy, ny) + abs(cy-ny)*0.22
            lh = abs(cy - ny) * 0.56
            if lw < 0.05: lw = 0.05
            if lh < 0.05: lh = 0.05
            rect(s, lx - 0.01, ly - 0.01, max(lw, 0.02), max(lh, 0.02), fill=TEAL)
        # Center "DNS" circle
        oval(s, cx-0.55, cy-0.55, 1.1, 1.1, fill=ORANGE)
        txt(s, "DNS", cx-0.55, cy-0.28, 1.1, 0.56, size=18, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
    return _section_slide(prs, "03", "DNS in the Wild",
                          "Real-world uses — email, services, and beyond",
                          13, icon_fn=icon)

def slide_15_section_security(prs):
    def icon(s):
        draw_lock(s, 9.2, 1.8, 3.2, 4.2, color=TEAL)
        # Exclamation badge
        oval(s, 11.8, 1.6, 0.8, 0.8, fill=ORANGE)
        txt(s, "!", 11.8, 1.62, 0.8, 0.78, size=28, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
    return _section_slide(prs, "04", "DNS Security",
                          "Protecting the lookup chain — DNSSEC, DoH, DoT",
                          15, icon_fn=icon)


# ── Slide 02 — Agenda ─────────────────────────────────────────────────────────
def slide_02_agenda(prs):
    s = blank_slide(prs)
    set_bg(s)
    title_bar(s, "What We'll Cover Today", 2)

    items = [
        ("01", "What is DNS?",         "The basics — names, IPs, and why we need it"),
        ("02", "How DNS Works",         "Recursive resolvers, root servers, authoritative DNS"),
        ("03", "Forward & Reverse DNS", "A/AAAA records vs PTR lookups"),
        ("04", "DNS Record Types",      "SOA, NS, A, AAAA, CNAME, MX, TXT, SRV, PTR, CAA, and more"),
        ("05", "DNS in the Wild",       "Email routing, service discovery, security, anti-spam"),
        ("06", "DNS Security",          "DNSSEC, DoH, DoT — protecting the lookup chain"),
    ]
    col_w = 4.1
    row_h = 1.6
    top = 1.6
    for i, (num, title, desc) in enumerate(items):
        col = i % 3
        row = i // 3
        x = 0.45 + col * (col_w + 0.25)
        y = top + row * (row_h + 0.18)
        rect(s, x, y, col_w, row_h, fill=BOX_MID, line_color=TEAL, line_w=Pt(1))
        # Orange left accent bar
        rect(s, x, y, 0.18, row_h, fill=ORANGE)
        txt(s, num,   x + 0.3, y + 0.12, 0.65, 0.55, size=26, bold=True, color=ORANGE)
        txt(s, title, x + 0.3, y + 0.62, col_w - 0.42, 0.42, size=16, bold=True, color=WHITE)
        txt(s, desc,  x + 0.3, y + 1.02, col_w - 0.42, 0.52, size=11, color=GRAY)
    return s


# ── Slide 04 — Problem DNS Solves ────────────────────────────────────────────
def slide_04_problem(prs):
    s = blank_slide(prs)
    set_bg(s)
    title_bar(s, "The Problem DNS Solves", 4)

    # Without DNS box
    rect(s, 0.4, 1.55, 5.8, 5.3, fill=BOX_DARK, line_color=ORANGE, line_w=Pt(1.5))
    rect(s, 0.4, 1.55, 5.8, 0.62, fill=BOX_ORANGE)
    txt(s, "✕  Without DNS", 0.65, 1.65, 5.2, 0.45, size=18, bold=True, color=ORANGE)
    # IP address visual
    rect(s, 0.7, 2.28, 5.1, 0.52, fill=BOX_MID, line_color=ORANGE, line_w=Pt(0.75))
    txt(s, "  142.250.80.46  (memorize this!)",
        0.78, 2.32, 4.9, 0.4, size=15, bold=True, color=ORANGE, font="Courier New")
    without = [
        "Every website needs an IP address",
        "IPs change when servers move",
        "No meaningful names — just numbers",
        "Sharing addresses is error-prone",
    ]
    for i, line in enumerate(without):
        txt(s, f"—  {line}", 0.75, 2.95 + i * 0.75, 5.1, 0.65, size=14, color=GRAY)

    # Big arrow
    rarrow(s, 6.3, 3.6, 0.75, 0.65)

    # With DNS box
    rect(s, 7.3, 1.55, 5.6, 5.3, fill=BOX_DARK, line_color=TEAL, line_w=Pt(1.5))
    rect(s, 7.3, 1.55, 5.6, 0.62, fill=BOX_MID)
    txt(s, "✓  With DNS", 7.55, 1.65, 5.0, 0.45, size=18, bold=True, color=TEAL)
    # Domain name visual
    rect(s, 7.6, 2.28, 4.9, 0.52, fill=BOX_MID, line_color=TEAL, line_w=Pt(0.75))
    txt(s, "  google.com  (easy!)",
        7.68, 2.32, 4.7, 0.4, size=15, bold=True, color=TEAL, font="Courier New")
    with_dns = [
        "DNS translates name → IP silently",
        "Servers can change IPs, names stay",
        "Human-readable, memorable names",
        "Hierarchical and globally distributed",
    ]
    for i, line in enumerate(with_dns):
        txt(s, f"✓  {line}", 7.55, 2.95 + i * 0.75, 5.1, 0.65, size=14, color=GRAY)

    # Bottom note
    rect(s, 0.4, 6.95, 12.5, 0.42, fill=BOX_MID)
    txt(s, "DNS: translating 142.250.80.46 into google.com since 1983  (RFC 882/883 → RFC 1034/1035)",
        0.6, 6.98, 12.1, 0.36, size=12, color=TEAL, italic=True)
    return s


# ── Slide 05 — DNS Lookup Flow ────────────────────────────────────────────────
def slide_05_how_lookup(prs):
    s = blank_slide(prs)
    set_bg(s)
    title_bar(s, "How a DNS Lookup Works", 5)

    steps = [
        ("1", "Your\nBrowser",         "You type\ncascadesteam.org"),
        ("2", "Recursive\nResolver",    "ISP or\n8.8.8.8"),
        ("3", "Root Name\nServer",      "Directs to\n.org TLD"),
        ("4", ".org TLD\nServer",       "Sends to\nauthoritative"),
        ("5", "Authoritative\nDNS",     "Returns IP\naddress!"),
    ]

    # Pentagon for step 1, Chevron for steps 2-5
    flow_w  = 2.42   # shape width
    overlap = 0.18   # chevron overlap
    step_x  = flow_w - overlap
    start_x = 0.28
    flow_y  = 1.55
    flow_h  = 2.35

    # Server icon labels above shapes
    server_icons = ["🖥️", "🔄", "🌍", "📂", "✅"]

    for i, (num, title, desc) in enumerate(steps):
        x = start_x + i * step_x
        stype = SH_PENTAGON if i == 0 else SH_CHEVRON
        fill_c = ORANGE if i == 4 else BOX_MID
        border_c = WHITE if i == 4 else TEAL

        sp = shape(s, stype, x, flow_y, flow_w, flow_h,
                   fill=fill_c, line_color=border_c, line_w=Pt(1.5))

        # Icon above shape
        txt(s, server_icons[i], x + 0.55, flow_y - 0.72, 1.3, 0.65,
            size=28, align=PP_ALIGN.CENTER)

        # Step number badge (circle approximated with oval)
        oval(s, x + 0.88 + (0.12 if i>0 else 0), flow_y + 0.18, 0.55, 0.55,
             fill=ORANGE if i < 4 else WHITE)
        txt(s, num,
            x + 0.88 + (0.12 if i>0 else 0), flow_y + 0.18, 0.55, 0.55,
            size=17, bold=True, align=PP_ALIGN.CENTER,
            color=WHITE if i < 4 else ORANGE)

        # Title text
        txt(s, title, x + 0.18 + (0.12 if i>0 else 0), flow_y + 0.82,
            flow_w - 0.42, 0.7, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Description
        txt(s, desc,  x + 0.18 + (0.12 if i>0 else 0), flow_y + 1.52,
            flow_w - 0.42, 0.75, size=11, color=GRAY, align=PP_ALIGN.CENTER)

    # Cascadesteam.org label below arrow flow
    txt(s, "cascadesteam.org  →  143.198.x.x",
        0.28, flow_y + flow_h + 0.1, 13.0, 0.38, size=13, bold=True,
        color=ORANGE, align=PP_ALIGN.CENTER, font="Courier New")

    # Caching note box
    rect(s, 0.4, 4.5, 12.55, 1.55, fill=BOX_DARK, line_color=TEAL, line_w=Pt(1))
    # Inner left accent
    rect(s, 0.4, 4.5, 0.1, 1.55, fill=TEAL)
    txt(s, "⚡  Caching & TTL", 0.65, 4.6, 3.5, 0.42, size=15, bold=True, color=TEAL)
    txt(s, "Answers are cached at each step for a TTL (Time To Live) period. "
           "If your resolver already has a cached answer, steps 3–5 are skipped entirely — "
           "making DNS extremely fast (sub-millisecond for cached queries). "
           "This is also why DNS changes take time to fully propagate globally.",
        0.65, 5.06, 12.1, 0.92, size=13, color=GRAY)
    return s


# ── Slide 07 — Address Records ────────────────────────────────────────────────
def slide_07_address_records(prs):
    s = blank_slide(prs)
    set_bg(s)
    title_bar(s, "Address Records: A, AAAA & CNAME", 7)

    records = [
        ("A",     "Address Record (IPv4)",
         "Maps a hostname to a 32-bit IPv4 address.\nMultiple A records = round-robin load balancing.",
         "cascadesteam.org.  A  143.198.56.78"),
        ("AAAA",  "IPv6 Address Record",
         "Maps a hostname to a 128-bit IPv6 address.\nModern dual-stack systems have both A and AAAA.",
         "cascadesteam.org.  AAAA  2600:1f18::1"),
        ("CNAME", "Canonical Name Record",
         "Creates an alias pointing to another hostname.\nCannot be used at the zone apex (root domain).",
         "www.example.com.  CNAME  example.com."),
    ]
    col_w = 4.1
    for i, (label, title, desc, code) in enumerate(records):
        x = 0.4 + i * (col_w + 0.25)
        rect(s, x, 1.55, col_w, 5.3, fill=BOX_DARK, line_color=TEAL, line_w=Pt(1))
        # Label badge
        rect(s, x + 0.15, 1.7, 1.1, 0.58, fill=ORANGE)
        txt(s, label, x + 0.15, 1.7, 1.1, 0.58, size=18, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
        txt(s, title, x + 0.15, 2.35, col_w - 0.3, 0.5, size=13, bold=True, color=WHITE)
        txt(s, desc,  x + 0.15, 2.85, col_w - 0.3, 1.35, size=12, color=GRAY)
        # Mini diagram: name → IP arrow
        rect(s, x + 0.15, 4.28, 1.35, 0.42, fill=BOX_MID, line_color=TEAL, line_w=Pt(0.75))
        txt(s, "hostname", x + 0.15, 4.31, 1.35, 0.36, size=10, color=TEAL,
            align=PP_ALIGN.CENTER, font="Courier New")
        rarrow(s, x + 1.56, 4.32, 0.55, 0.38, fill=TEAL)
        rect(s, x + 2.18, 4.28, 1.72, 0.42, fill=BOX_MID, line_color=ORANGE, line_w=Pt(0.75))
        txt(s, "IP address", x + 2.18, 4.31, 1.72, 0.36, size=10, color=ORANGE,
            align=PP_ALIGN.CENTER, font="Courier New")
        # Code example
        code_box(s, code, x + 0.15, 4.82, col_w - 0.3, 0.62)
    return s


# ── Slide 08 — MX and TXT ────────────────────────────────────────────────────
def slide_08_mx_txt(prs):
    s = blank_slide(prs)
    set_bg(s)
    title_bar(s, "Mail & Text Records: MX and TXT", 8)

    # MX panel (left half)
    rect(s, 0.4, 1.55, 5.8, 5.3, fill=BOX_DARK, line_color=ORANGE, line_w=Pt(1.5))
    rect(s, 0.55, 1.7, 0.9, 0.55, fill=ORANGE)
    txt(s, "MX", 0.55, 1.7, 0.9, 0.55, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "Mail Exchanger Record", 1.55, 1.78, 4.4, 0.45, size=15, bold=True, color=WHITE)
    txt(s, "Tells the world which mail servers accept email for your domain.\n"
           "Lower priority number = higher priority. Backup servers take over if primary fails.",
        0.6, 2.32, 5.4, 1.05, size=12, color=GRAY)
    # MX priority visual
    for i, (prio, srv, bar_w) in enumerate([("10", "mail1 (primary)", 4.5),
                                              ("20", "mail2 (backup)",  3.5),
                                              ("30", "mail3 (tertiary)",2.5)]):
        y = 3.45 + i * 0.52
        rect(s, 0.6, y, bar_w, 0.38, fill=BOX_MID, line_color=TEAL, line_w=Pt(0.5))
        rect(s, 0.6, y, 0.5, 0.38, fill=ORANGE)
        txt(s, prio, 0.6, y + 0.03, 0.5, 0.32, size=11, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
        txt(s, srv, 1.15, y + 0.05, 3.5, 0.3, size=11, color=GRAY)
    txt(s, "Google Workspace, Office 365, and Proton Mail all require\nMX records to route email to their servers.",
        0.6, 5.14, 5.4, 0.65, size=11, color=GRAY, italic=True)

    # TXT panel (right half)
    rect(s, 6.5, 1.55, 6.45, 5.3, fill=BOX_DARK, line_color=TEAL, line_w=Pt(1.5))
    rect(s, 6.65, 1.7, 0.9, 0.55, fill=TEAL)
    txt(s, "TXT", 6.65, 1.7, 0.9, 0.55, size=18, bold=True,
        color=RGBColor(0x06, 0x1d, 0x29), align=PP_ALIGN.CENTER)
    txt(s, "Text Record — Swiss Army Knife", 7.65, 1.78, 5.1, 0.45, size=15, bold=True, color=WHITE)
    txt(s, "TXT records store arbitrary text. Critical for domain verification and email security:",
        6.65, 2.32, 6.1, 0.65, size=12, color=GRAY)

    sub_items = [
        ("SPF",    "Specifies which servers may send email",        "v=spf1 include:_spf.google.com ~all"),
        ("DKIM",   "Public key to verify email signatures",         "v=DKIM1; k=rsa; p=MIGfMA0..."),
        ("DMARC",  "Email auth policy (p=reject/quarantine)",       "v=DMARC1; p=reject; rua=mailto:..."),
        ("Verify", "Domain ownership for external services",        "google-site-verification=abc123"),
    ]
    for i, (label, desc, code) in enumerate(sub_items):
        y = 3.08 + i * 0.72
        rect(s, 6.65, y, 0.8, 0.32, fill=ORANGE)
        txt(s, label, 6.65, y, 0.8, 0.32, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, desc, 7.52, y, 3.0, 0.32, size=11, color=GRAY)
        txt(s, code, 10.6, y, 2.2, 0.32, size=10, color=TEAL, font="Courier New")
    return s


# ── Slide 09 — SOA and NS ─────────────────────────────────────────────────────
def slide_09_soa_ns(prs):
    s = blank_slide(prs)
    set_bg(s)
    title_bar(s, "Zone Control Records: SOA and NS", 9)

    # SOA panel
    rect(s, 0.4, 1.55, 5.9, 4.3, fill=BOX_DARK, line_color=ORANGE, line_w=Pt(1.5))
    rect(s, 0.55, 1.7, 0.9, 0.55, fill=ORANGE)
    txt(s, "SOA", 0.55, 1.7, 0.9, 0.55, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "Start of Authority — Every Zone Has Exactly One",
        1.55, 1.78, 4.55, 0.45, size=13, bold=True, color=WHITE)
    txt(s, "The boss record of every DNS zone. Contains: primary name server, "
           "admin email, serial number (version counter), refresh/retry/expire timers, "
           "and negative TTL.",
        0.6, 2.32, 5.5, 1.1, size=12, color=GRAY)
    code_box(s, "cascadesteam.org.  SOA\n  ns1.hover.com.  dns.hover.com.\n  2024040201  ; serial\n  3600        ; refresh\n  900         ; retry\n  604800      ; expire\n  300         ; negative TTL",
             0.6, 3.5, 5.5, 1.8)

    # NS panel
    rect(s, 6.65, 1.55, 6.3, 4.3, fill=BOX_DARK, line_color=TEAL, line_w=Pt(1.5))
    rect(s, 6.8, 1.7, 0.75, 0.55, fill=TEAL)
    txt(s, "NS", 6.8, 1.7, 0.75, 0.55, size=18, bold=True,
        color=RGBColor(0x06, 0x1d, 0x29), align=PP_ALIGN.CENTER)
    txt(s, "Name Server — Who Answers Your Domain",
        7.65, 1.78, 5.1, 0.45, size=13, bold=True, color=WHITE)
    txt(s, "NS records delegate authority for a zone to specific name servers. "
           "Always have at least 2 for redundancy. Also delegates subdomains to "
           "different DNS providers.",
        6.8, 2.32, 6.0, 1.1, size=12, color=GRAY)
    code_box(s, "cascadesteam.org.  NS  ns1.hover.com.\ncascadesteam.org.  NS  ns2.hover.com.\n\n; Subdomain delegation:\ndevops.example.com.  NS  ns1.internal.",
             6.8, 3.5, 6.0, 1.6)

    # DNS Hierarchy tree at bottom
    rect(s, 0.4, 6.08, 12.55, 1.28, fill=BOX_MID, line_color=TEAL, line_w=Pt(0.75))
    txt(s, "DNS Delegation Chain:", 0.65, 6.15, 3.0, 0.35, size=12, bold=True, color=TEAL)
    # Chain: ROOT → .ORG → cascadesteam.org → ns1/ns2.hover.com
    chain = [
        (". (Root)", ORANGE),
        (".org TLD", TEAL),
        ("cascadesteam.org", BOX_DARK),
        ("ns1.hover.com\nns2.hover.com", BOX_DARK),
    ]
    bw = 2.6
    for i, (label, color) in enumerate(chain):
        bx = 0.55 + i * (bw + 0.05)
        rect(s, bx, 6.52, bw - 0.05, 0.72, fill=color, line_color=TEAL, line_w=Pt(0.75))
        txt(s, label, bx + 0.08, 6.55, bw - 0.2, 0.62, size=11, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        if i < 3:
            rarrow(s, bx + bw, 6.67, 0.1, 0.4, fill=ORANGE)
    return s


# ── Slide 10 — PTR Records ────────────────────────────────────────────────────
def slide_10_ptr(prs):
    s = blank_slide(prs)
    set_bg(s)
    title_bar(s, "Reverse DNS: PTR Records", 10)

    rect(s, 0.4, 1.55, 12.5, 1.35, fill=BOX_DARK, line_color=TEAL, line_w=Pt(1))
    rect(s, 0.55, 1.68, 0.85, 0.55, fill=ORANGE)
    txt(s, "PTR", 0.55, 1.68, 0.85, 0.55, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "Pointer Record — DNS in Reverse.  Maps IP → hostname in in-addr.arpa (IPv4) "
           "or ip6.arpa (IPv6), controlled by the IP owner (ISP or hosting provider).",
        1.5, 1.7, 11.2, 0.55, size=13, color=GRAY)

    # Visual forward/reverse diagram
    # Forward: hostname box → arrow → IP box
    rect(s, 0.4, 3.08, 3.5, 0.72, fill=BOX_MID, line_color=TEAL, line_w=Pt(1))
    txt(s, "mail.example.com", 0.5, 3.14, 3.3, 0.58, size=14, bold=True,
        color=TEAL, align=PP_ALIGN.CENTER, font="Courier New")

    rarrow(s, 4.0, 3.18, 1.1, 0.52, fill=TEAL)
    txt(s, "A record", 4.0, 3.1, 1.1, 0.3, size=10, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

    rect(s, 5.25, 3.08, 3.5, 0.72, fill=BOX_MID, line_color=ORANGE, line_w=Pt(1))
    txt(s, "203.0.113.42", 5.35, 3.14, 3.3, 0.58, size=14, bold=True,
        color=ORANGE, align=PP_ALIGN.CENTER, font="Courier New")

    # Reverse arrow (left-pointing)
    shape(s, 34, 4.0, 3.88, 1.1, 0.52, fill=ORANGE)  # Left arrow = 34
    txt(s, "PTR record", 4.0, 3.82, 1.1, 0.3, size=10, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

    rect(s, 0.4, 3.88, 3.5, 0.72, fill=BOX_DARK, line_color=ORANGE, line_w=Pt(1))
    txt(s, "42.113.0.203\n.in-addr.arpa", 0.5, 3.9, 3.3, 0.68, size=12, bold=True,
        color=ORANGE, align=PP_ALIGN.CENTER, font="Courier New")

    txt(s, "Lookup direction", 5.55, 3.88, 3.1, 0.72, size=11, color=GRAY,
        align=PP_ALIGN.CENTER, italic=True)

    txt(s, "Why PTR Records Matter:", 0.4, 4.82, 8, 0.38, size=15, bold=True, color=WHITE)
    reasons = [
        ("Email Deliverability", "Mail servers check PTR of sending IP. No match = rejected or spam."),
        ("Network Troubleshooting", "traceroute and ping show hostnames, not raw IPs."),
        ("Log Readability", "Server logs show 'mail.example.com' not '203.0.113.42'."),
        ("Security Verification", "IDS and financial systems require matching forward+reverse DNS."),
    ]
    for i, (label, desc) in enumerate(reasons):
        col = i % 2
        row = i // 2
        x = 0.4 + col * 6.3
        y = 5.28 + row * 1.05
        rect(s, x, y, 6.05, 0.9, fill=BOX_MID, line_color=TEAL, line_w=Pt(0.75))
        rect(s, x, y, 0.12, 0.9, fill=ORANGE)
        txt(s, label, x + 0.25, y + 0.08, 5.65, 0.35, size=13, bold=True, color=WHITE)
        txt(s, desc,  x + 0.25, y + 0.44, 5.65, 0.42, size=11, color=GRAY)
    return s


# ── Slide 11 — SRV Records ────────────────────────────────────────────────────
def slide_11_srv(prs):
    s = blank_slide(prs)
    set_bg(s)
    title_bar(s, "Service Discovery: SRV Records", 11)

    rect(s, 0.4, 1.55, 12.5, 1.4, fill=BOX_DARK, line_color=TEAL, line_w=Pt(1))
    rect(s, 0.55, 1.7, 0.85, 0.55, fill=ORANGE)
    txt(s, "SRV", 0.55, 1.7, 0.85, 0.55, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "Service Locator Record — tells clients where to find a service (hostname AND port). "
           "Applications auto-configure from DNS — no hardcoded IPs or ports needed.",
        1.5, 1.72, 11.2, 0.55, size=13, color=GRAY)
    txt(s, "Format:  _service._protocol.name.  TTL  IN  SRV  priority  weight  port  target",
        0.6, 2.28, 12.1, 0.38, size=12, color=TEAL, font="Courier New")
    code_box(s, "_xmpp-client._tcp.jabber.org.  86400  IN  SRV  5  50  5222  xmpp1.jabber.org.",
             0.6, 2.7, 12.1, 0.52)

    fields = [
        ("_service",  "Protocol name (e.g., _xmpp, _sip, _ldap, _kerberos)"),
        ("_protocol", "_tcp or _udp"),
        ("priority",  "Lower = tried first (like MX). Used for failover."),
        ("weight",    "Load balancing between equal-priority records"),
        ("port",      "The TCP/UDP port the service listens on"),
        ("target",    "The hostname of the server providing the service"),
    ]
    col_w = 6.1
    for i, (label, desc) in enumerate(fields):
        col = i % 2
        row = i // 2
        x = 0.4 + col * (col_w + 0.2)
        y = 3.35 + row * 0.72
        rect(s, x, y, col_w, 0.65, fill=BOX_MID, line_color=TEAL, line_w=Pt(0.5))
        rect(s, x, y, 1.5, 0.65, fill=BOX_DARK)
        txt(s, label,  x + 0.1, y + 0.12, 1.3, 0.42, size=12, bold=True, color=ORANGE, font="Courier New")
        txt(s, desc,   x + 1.6, y + 0.12, col_w - 1.7, 0.42, size=12, color=GRAY)

    txt(s, "Used by: XMPP, SIP/VoIP, LDAP, Kerberos, Active Directory, CalDAV/CardDAV, Minecraft, Kubernetes",
        0.4, 5.62, 12.5, 0.4, size=12, color=ORANGE, italic=True)
    return s


# ── Slide 12 — More Record Types ─────────────────────────────────────────────
def slide_12_more_records(prs):
    s = blank_slide(prs)
    set_bg(s)
    title_bar(s, "More Record Types Worth Knowing", 12)

    records = [
        ("CAA",   "CA Authorization",
         "Specifies which CAs may issue TLS certificates for your domain.",
         "example.com.  CAA  0 issue \"letsencrypt.org\""),
        ("NAPTR", "Naming Authority Pointer",
         "Used in VoIP/SIP and ENUM to translate phone numbers to URIs.",
         "example.com.  NAPTR  10 10 \"u\" \"E2U+sip\" ..."),
        ("TLSA",  "TLS Authentication (DANE)",
         "Pins a TLS certificate to a port/protocol via DNS.",
         "_443._tcp.example.com.  TLSA  3 1 1 abc123..."),
        ("HINFO", "Host Information",
         "Stores CPU and OS type. Rarely used today — historical curiosity.",
         "host.example.com.  HINFO  \"x86-64\" \"Linux\""),
        ("LOC",   "Geographic Location",
         "Stores GPS coordinates for a host. Standardized in RFC 1876.",
         "host.example.com.  LOC  48 44 N 122 28 W 60m"),
        ("SSHFP", "SSH Fingerprint",
         "Stores SSH host key fingerprints. With DNSSEC, clients auto-verify.",
         "host.example.com.  SSHFP  2 1 7491973e..."),
    ]
    col_w = 4.1
    col_h = 2.5
    for i, (label, title, desc, code) in enumerate(records):
        col = i % 3
        row = i // 3
        x = 0.4 + col * (col_w + 0.22)
        y = 1.55 + row * (col_h + 0.22)
        rect(s, x, y, col_w, col_h, fill=BOX_DARK, line_color=TEAL, line_w=Pt(1))
        rect(s, x + 0.12, y + 0.12, 1.0, 0.45, fill=ORANGE)
        txt(s, label,  x + 0.12, y + 0.12, 1.0, 0.45, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, title,  x + 1.2,  y + 0.16, col_w - 1.32, 0.42, size=12, bold=True, color=WHITE)
        txt(s, desc,   x + 0.12, y + 0.65, col_w - 0.24, 0.82, size=11, color=GRAY)
        code_box(s, code, x + 0.12, y + 1.52, col_w - 0.24, 0.55)
    return s


# ── Slide 14 — Email Security ─────────────────────────────────────────────────
def slide_14_email_security(prs):
    s = blank_slide(prs)
    set_bg(s)
    title_bar(s, "Email Anti-Spam Trio: SPF + DKIM + DMARC", 14)

    # Email journey flow strip
    rect(s, 0.4, 1.52, 12.55, 0.7, fill=BOX_DARK, line_color=TEAL, line_w=Pt(0.75))
    flow_items = [
        ("📧 SEND",       BOX_MID,   WHITE),
        ("SPF?",          ORANGE,    WHITE),
        ("DKIM?",         TEAL,      BOX_DARK),
        ("DMARC?",        BOX_MID,   WHITE),
        ("✓ DELIVER",     GREEN,     WHITE),
        ("✗ BLOCK",       RED_CLR,   WHITE),
    ]
    fw = 1.68
    for i, (label, bg, fg) in enumerate(flow_items):
        bx = 0.52 + i * (fw + 0.12)
        if i == 4:
            bx = 0.52 + i * (fw + 0.12) + 0.1
        if i == 5:
            bx = 0.52 + 4 * (fw + 0.12) + 0.1 + fw + 0.08
        rect(s, bx, 1.6, fw, 0.54, fill=bg, line_color=TEAL, line_w=Pt(0.5))
        txt(s, label, bx + 0.06, 1.64, fw - 0.12, 0.44, size=13, bold=True,
            color=fg, align=PP_ALIGN.CENTER)
        if i < 3:
            rarrow(s, bx + fw, 1.73, 0.15, 0.3, fill=GRAY)

    # Branch label
    txt(s, "pass all → deliver  |  fail any → quarantine/reject",
        9.1, 1.6, 3.6, 0.54, size=9, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

    panels = [
        ("SPF",   "Sender Policy Framework",   ORANGE,
         "A TXT record listing all servers authorized to send email from your domain. "
         "Receiving servers check if the sender's IP is on the list.",
         "v=spf1 include:_spf.google.com ip4:203.0.113.0/24 ~all",
         "Soft fail (~all) or hard fail (-all) if IP isn't listed"),
        ("DKIM",  "DomainKeys Identified Mail", TEAL,
         "Your mail server signs outgoing emails with a private key. "
         "The public key is a TXT record. Recipients verify the signature.",
         "selector1._domainkey.example.com. TXT \"v=DKIM1; k=rsa; p=MIGf...\"",
         "Signature mismatch = email was altered in transit"),
        ("DMARC", "Domain-based Message Auth",  RGBColor(0x7a, 0xcc, 0x44),
         "Builds on SPF and DKIM. Tells receivers what to do with failing messages "
         "(none/quarantine/reject) and sends reports back to you.",
         "_dmarc.example.com. TXT \"v=DMARC1; p=reject; rua=mailto:dmarc@...\"",
         "p=reject = unauthenticated email is dropped entirely"),
    ]
    col_w = 4.1
    for i, (label, title, accent, desc, code, note) in enumerate(panels):
        x = 0.4 + i * (col_w + 0.22)
        rect(s, x, 2.32, col_w, 4.9, fill=BOX_DARK, line_color=accent, line_w=Pt(1.5))
        rect(s, x + 0.12, 2.45, 1.1, 0.55, fill=accent)
        label_color = WHITE if accent != TEAL else RGBColor(0x06, 0x1d, 0x29)
        txt(s, label, x + 0.12, 2.45, 1.1, 0.55, size=17, bold=True, color=label_color, align=PP_ALIGN.CENTER)
        txt(s, title, x + 1.32, 2.52, col_w - 1.44, 0.46, size=13, bold=True, color=WHITE)
        txt(s, "How it works:", x + 0.15, 3.1, col_w - 0.3, 0.3, size=11, bold=True, color=accent)
        txt(s, desc,   x + 0.15, 3.42, col_w - 0.3, 1.35, size=11, color=GRAY)
        code_box(s, code, x + 0.15, 4.88, col_w - 0.3, 0.55)
        txt(s, note,   x + 0.15, 5.5, col_w - 0.3, 0.62, size=10, color=GRAY, italic=True)
    return s


# ── Slide 16 — DNSSEC / DoH / DoT ────────────────────────────────────────────
def slide_16_dns_security(prs):
    s = blank_slide(prs)
    set_bg(s)
    title_bar(s, "DNS Security: DNSSEC, DoH, and DoT", 16)

    # Attack scenario visual banner
    rect(s, 0.4, 1.52, 12.55, 0.72, fill=RGBColor(0x3a, 0x15, 0x05),
         line_color=ORANGE, line_w=Pt(1))
    txt(s, "⚠  Classic DNS is unencrypted and unauthenticated — attackers can intercept, "
           "modify, or spoof responses pointing users to malicious servers.",
        0.6, 1.58, 12.1, 0.58, size=12, color=WHITE)

    # Before/After comparison strip
    rect(s, 0.4, 2.34, 6.1, 0.58, fill=RGBColor(0x2a, 0x07, 0x07), line_color=RED_CLR, line_w=Pt(0.75))
    txt(s, "🔓  BEFORE:  DNS query (UDP:53) — visible to anyone on the network",
        0.6, 2.4, 5.8, 0.45, size=11, color=RED_CLR)
    rect(s, 6.6, 2.34, 6.35, 0.58, fill=RGBColor(0x07, 0x2a, 0x12), line_color=GREEN, line_w=Pt(0.75))
    txt(s, "🔒  AFTER:   Encrypted tunnel (port 443/853) — only resolver sees queries",
        6.8, 2.4, 6.05, 0.45, size=11, color=GREEN)

    panels = [
        ("DNSSEC", "DNS Security Extensions", ORANGE, [
            "Adds digital signatures to DNS records",
            "Chain of trust: root → TLD → zone",
            "Validates responses haven't been tampered with",
            "Does NOT encrypt traffic — just signs it",
            "Special records: RRSIG, DNSKEY, DS, NSEC",
        ]),
        ("DoH", "DNS over HTTPS  (RFC 8484)", TEAL, [
            "Wraps DNS queries inside HTTPS (port 443)",
            "Looks like normal web traffic — hard to block",
            "Encrypted end-to-end between client and resolver",
            "Supported by Firefox, Chrome, Windows 11",
            "Providers: Cloudflare 1.1.1.1, Google 8.8.8.8",
        ]),
        ("DoT", "DNS over TLS  (RFC 7858)", GREEN, [
            "Wraps DNS queries inside TLS (port 853)",
            "Dedicated port — easier to monitor",
            "Full encryption between client and resolver",
            "Network admins can see and block DoT traffic",
            "Android 9+ uses DoT as 'Private DNS' feature",
        ]),
    ]
    col_w = 4.1
    for i, (label, title, accent, bullets) in enumerate(panels):
        x = 0.4 + i * (col_w + 0.22)
        rect(s, x, 3.05, col_w, 4.0, fill=BOX_DARK, line_color=accent, line_w=Pt(1.5))
        # Top label badge
        rect(s, x + 0.12, 3.18, 1.45, 0.58, fill=accent)
        label_color = WHITE if accent != TEAL else RGBColor(0x06, 0x1d, 0x29)
        txt(s, label, x + 0.12, 3.18, 1.45, 0.58, size=18, bold=True, color=label_color, align=PP_ALIGN.CENTER)
        txt(s, title, x + 1.65, 3.25, col_w - 1.77, 0.46, size=12, bold=True, color=WHITE)
        # Lock icon for DoH/DoT
        if i > 0:
            draw_lock(s, x + col_w - 0.95, 3.12, 0.7, 0.88,
                      color=accent if i == 1 else GREEN)
        for j, bullet in enumerate(bullets):
            txt(s, f"•  {bullet}", x + 0.18, 3.9 + j * 0.62, col_w - 0.3, 0.55, size=11, color=GRAY)
    return s


# ── Slide 17 — Quick Reference Table ─────────────────────────────────────────
def slide_17_reference(prs):
    s = blank_slide(prs)
    set_bg(s)
    title_bar(s, "DNS Record Quick Reference", 17)

    rows = [
        ("A",     "Name → IPv4",        "Most fundamental record",     "203.0.113.42"),
        ("AAAA",  "Name → IPv6",        "128-bit, dual-stack",         "2001:db8::1"),
        ("CNAME", "Alias → Name",       "Cannot be at zone apex",      "www → example.com."),
        ("MX",    "Email routing",      "Priority (lower = first)",    "10 mail.example.com."),
        ("TXT",   "Arbitrary text",     "SPF, DKIM, DMARC, verify",    "v=spf1 include:... ~all"),
        ("SOA",   "Zone authority",     "One per zone, serial #",      "ns1. admin. 20240401 ..."),
        ("NS",    "Authoritative svrs", "Delegation, min 2 rec.",      "ns1.hover.com."),
        ("PTR",   "IP → Name (reverse)","Owned by IP block owner",     "42.113.0.203.in-addr.arpa"),
        ("SRV",   "Service discovery",  "Priority, weight, port, target","_xmpp._tcp 5 50 5222 host."),
        ("CAA",   "CA authorization",   "Which CA may issue SSL certs","0 issue \"letsencrypt.org\""),
        ("SSHFP", "SSH key fingerprint","Validated by DNSSEC",         "2 1 7491973e5f8b39..."),
        ("TLSA",  "TLS cert pin",       "DANE — cert via DNS",         "3 1 1 abc123def456..."),
    ]

    cols = [("Record", 1.15), ("Purpose", 2.5), ("Key Info", 3.5), ("Example Value", 5.0)]
    header_y = 1.52
    x = 0.35
    for col_title, col_w in cols:
        rect(s, x, header_y, col_w - 0.05, 0.5, fill=ORANGE)
        txt(s, col_title, x + 0.1, header_y + 0.05, col_w - 0.25, 0.4,
            size=13, bold=True, color=WHITE)
        x += col_w

    row_h = 0.44
    for ri, (rec, purpose, key, example) in enumerate(rows):
        y = 2.08 + ri * row_h
        bg = BOX_MID if ri % 2 == 0 else BOX_DARK
        x = 0.35
        widths = [1.15, 2.5, 3.5, 5.0]
        vals   = [rec, purpose, key, example]
        colors = [ORANGE, WHITE, GRAY, TEAL]
        fonts  = ["Calibri", "Calibri", "Calibri", "Courier New"]
        bolds  = [True, False, False, False]
        for ci, (v, w, c, f, b) in enumerate(zip(vals, widths, colors, fonts, bolds)):
            rect(s, x, y, w - 0.05, row_h - 0.02, fill=bg, line_color=BG, line_w=Pt(0.5))
            txt(s, v, x + 0.1, y + 0.04, w - 0.25, row_h - 0.08,
                size=11, color=c, font=f, bold=b)
            x += w
    return s


# ── Slide 18 — Closing ────────────────────────────────────────────────────────
def slide_18_closing(prs):
    s = blank_slide(prs)
    set_bg(s)
    topo_lines(s)

    rect(s, 0, 0, 0.12, 7.5, fill=TEAL)

    txt(s, "Key Takeaways", 0.5, 0.25, 9, 0.7, size=36, bold=True, color=WHITE)
    rect(s, 0.5, 0.92, 8.5, 0.04, fill=ORANGE)

    takeaways = [
        ("01", "DNS is infrastructure",
         "Every internet connection depends on DNS. Foundational to networking, security, and dev work."),
        ("02", "Records have specific jobs",
         "A/AAAA serve IPs, MX routes email, SRV discovers services, TXT handles auth, PTR enables reverse."),
        ("03", "DNS powers email security",
         "SPF, DKIM, and DMARC are DNS-based. Proper config stops spoofing and protects your domain."),
        ("04", "DNS is increasingly secure",
         "DNSSEC signs records, DoH/DoT encrypt queries. The old plaintext UDP era is ending."),
        ("05", "TTL controls propagation",
         "DNS changes don't happen instantly. TTL governs how long caches hold answers — plan accordingly."),
    ]
    for i, (num, title, desc) in enumerate(takeaways):
        y = 1.1 + i * 1.1
        rect(s, 0.5, y, 8.5, 0.98, fill=BOX_MID, line_color=TEAL, line_w=Pt(0.75))
        rect(s, 0.5, y, 0.68, 0.98, fill=ORANGE)
        txt(s, num, 0.5, y + 0.2, 0.68, 0.55, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, title, 1.3, y + 0.08, 7.5, 0.38, size=15, bold=True, color=WHITE)
        txt(s, desc,  1.3, y + 0.46, 7.5, 0.46, size=12, color=GRAY)

    # "Go Further" sidebar
    rect(s, 9.3, 1.0, 3.8, 5.7, fill=BOX_DARK, line_color=ORANGE, line_w=Pt(1.5))
    rect(s, 9.3, 1.0, 3.8, 0.5, fill=ORANGE)
    txt(s, "Go Further", 9.5, 1.06, 3.4, 0.38, size=16, bold=True, color=WHITE)
    resources = [
        "RFC 1034/1035 — DNS spec",
        "RFC 1912 — Common DNS errors",
        "IANA DNS Parameters",
        "MXToolbox.com",
        "dnsviz.net — DNSSEC viz",
        "dig, nslookup, drill",
        "Pi-hole for DNS filtering",
        "Cloudflare Learning DNS",
    ]
    for i, r in enumerate(resources):
        txt(s, f"• {r}", 9.5, 1.58 + i * 0.58, 3.4, 0.5, size=11, color=GRAY)

    txt(s, "Questions?", 0.5, 6.5, 5, 0.65, size=30, bold=True, color=TEAL)
    logo_img(s, 10.2, 6.3, w=2.85)
    footer(s, 18)
    return s


# ── Main Build ────────────────────────────────────────────────────────────────

def build_pptx():
    prs = new_prs()

    builders = [
        slide_01_title,
        slide_02_agenda,
        slide_03_section_what_is_dns,
        slide_04_problem,
        slide_05_how_lookup,
        slide_06_section_records,
        slide_07_address_records,
        slide_08_mx_txt,
        slide_09_soa_ns,
        slide_10_ptr,
        slide_11_srv,
        slide_12_more_records,
        slide_13_section_wild,
        slide_14_email_security,
        slide_15_section_security,
        slide_16_dns_security,
        slide_17_reference,
        slide_18_closing,
    ]

    for fn in builders:
        print(f"  Building {fn.__name__}...")
        fn(prs)

    prs.save(PPTX_OUT)
    print(f"\n✓ PPTX saved: {PPTX_OUT}")


# ── Markdown (Obsidian Slides) ─────────────────────────────────────────────────

MARKDOWN = '''---
tags: [open-source, dns, presentation, cascade-steam]
cssclasses: [slides]
---

<!-- Obsidian Slides — use the Slides core plugin to present -->

# DNS
## More Than Just a Lookup Table

Domain Name System — A deep dive into the internet\'s phone book, directory, security layer, and so much more.

*Cascade STEAM — Technology Education Series*

---

## What We\'ll Cover Today

| # | Topic | Summary |
|---|-------|---------|
| 01 | What is DNS? | The basics — names, IPs, and why we need it |
| 02 | How DNS Works | Recursive resolvers, root servers, authoritative DNS |
| 03 | Forward & Reverse DNS | A/AAAA records vs PTR lookups |
| 04 | DNS Record Types | SOA, NS, A, AAAA, CNAME, MX, TXT, SRV, PTR, CAA, and more |
| 05 | DNS in the Wild | Email routing, service discovery, security, anti-spam |
| 06 | DNS Security | DNSSEC, DoH, DoT — protecting the lookup chain |

---

# 01 — What is DNS?

*The basics — names, IPs, and the phone book of the internet*

---

## The Problem DNS Solves

**Without DNS**
- You\'d need to memorize: `142.250.80.46`
- Every website needs an IP address
- IPs change when servers move
- No meaningful names — just numbers

**With DNS**
- You type: `google.com`
- DNS translates name → IP silently
- Servers can change IPs, names stay the same
- Human-readable, memorable, hierarchical

> DNS: translating 142.250.80.46 into google.com since 1983 (RFC 882/883 → RFC 1034/1035)

---

## How a DNS Lookup Works

```
🖥️ Your Browser  →  🔄 Recursive Resolver  →  🌍 Root NS  →  📂 .org TLD  →  ✅ Authoritative DNS
```

1. **Your Browser** — You type `cascadesteam.org`
2. **Recursive Resolver** — Your ISP or `8.8.8.8` does the heavy lifting
3. **Root Name Server** — Directs to `.org` TLD servers
4. **.org TLD Server** — Directs to authoritative DNS for `cascadesteam.org`
5. **Authoritative DNS** — Returns the actual IP address!

`cascadesteam.org → 143.198.x.x`

> ⚡ **Caching & TTL:** Cached answers skip steps 3–5 — making DNS sub-millisecond for warm queries.

---

# 02 — DNS Record Types

*The Building Blocks — every record has a job*

---

## Address Records: A, AAAA & CNAME

### `A` — Address Record (IPv4)
`hostname → IP address`
Maps a hostname to a 32-bit IPv4 address. Multiple A records = round-robin load balancing.
```
cascadesteam.org.  A  143.198.56.78
```

### `AAAA` — IPv6 Address Record
Maps a hostname to a 128-bit IPv6 address. Modern dual-stack systems have both.
```
cascadesteam.org.  AAAA  2600:1f18::1
```

### `CNAME` — Canonical Name Record
Creates an alias pointing to another hostname. Cannot be used at the zone apex.
```
www.example.com.  CNAME  example.com.
```

---

## Mail & Text Records: MX and TXT

### `MX` — Mail Exchanger
Priority-based mail server routing. Lower number = higher priority.
```
example.com.  10  mail1.example.com.   ← primary
example.com.  20  mail2.example.com.   ← backup
example.com.  30  mail3.example.com.   ← tertiary
```

### `TXT` — Text Record (Swiss Army Knife)

| Type | Purpose | Example |
|------|---------|---------|
| SPF | Authorized sending servers | `v=spf1 include:_spf.google.com ~all` |
| DKIM | Public key for email signatures | `v=DKIM1; k=rsa; p=MIGfMA0...` |
| DMARC | Email authentication policy | `v=DMARC1; p=reject; rua=mailto:...` |
| Verify | Domain ownership | `google-site-verification=abc123` |

---

## Zone Control Records: SOA and NS

### `SOA` — Start of Authority
Every DNS zone has exactly one. Contains: primary NS, admin email, serial, timers.
```
cascadesteam.org.  SOA
  ns1.hover.com.  dns.hover.com.
  2024040201  ; serial
  3600        ; refresh
```

### `NS` — Name Server
Delegates authority for a zone to specific name servers. Min 2 for redundancy.
```
cascadesteam.org.  NS  ns1.hover.com.
cascadesteam.org.  NS  ns2.hover.com.
```

**DNS Delegation Chain:** `. (Root)` → `.org TLD` → `cascadesteam.org` → `ns1/ns2.hover.com`

---

## Reverse DNS: PTR Records

### `PTR` — Pointer Record

```
Forward (A):     mail.example.com  →  203.0.113.42
Reverse (PTR):   42.113.0.203.in-addr.arpa  →  mail.example.com
```

**Why PTR Records Matter:**
- **Email Deliverability** — No PTR = spam or rejection
- **Network Troubleshooting** — `traceroute` shows hostnames not raw IPs
- **Log Readability** — Human-readable server logs
- **Security Verification** — IDS and financial systems require matching forward+reverse

---

## Service Discovery: SRV Records

### `SRV` — Service Locator
Format: `_service._protocol.name.  TTL  IN  SRV  priority  weight  port  target`
```
_xmpp-client._tcp.jabber.org.  86400  IN  SRV  5  50  5222  xmpp1.jabber.org.
```

| Field | Meaning |
|-------|---------|
| `_service` | Protocol (_xmpp, _sip, _ldap, _kerberos) |
| `priority` | Lower = tried first (failover) |
| `weight` | Load balancing between equal-priority records |
| `port` | TCP/UDP port the service listens on |

*Used by: XMPP, SIP/VoIP, LDAP, Active Directory, Minecraft, Kubernetes*

---

## More Record Types Worth Knowing

| Record | Purpose | Example |
|--------|---------|---------|
| `CAA` | Which CAs may issue TLS certs | `0 issue "letsencrypt.org"` |
| `NAPTR` | VoIP phone number → URI | `10 10 "u" "E2U+sip" ...` |
| `TLSA` | TLS cert pin via DNS (DANE) | `3 1 1 abc123...` |
| `HINFO` | Host CPU/OS info (legacy) | `"x86-64" "Linux"` |
| `LOC` | GPS coordinates (RFC 1876) | `48 44 N 122 28 W 60m` |
| `SSHFP` | SSH host key fingerprint | `2 1 7491973e...` |

---

# 03 — DNS in the Wild

*Real-world uses — email, services, and beyond*

---

## Email Anti-Spam Trio: SPF + DKIM + DMARC

**Email Journey:** `📧 SEND` → `SPF?` → `DKIM?` → `DMARC?` → `✓ DELIVER` or `✗ BLOCK`

### SPF — Sender Policy Framework
```
v=spf1 include:_spf.google.com ip4:203.0.113.0/24 ~all
```
Lists all servers authorized to send from your domain. Soft fail `~all` or hard fail `-all`.

### DKIM — DomainKeys Identified Mail
```
selector1._domainkey.example.com.  TXT  "v=DKIM1; k=rsa; p=MIGfMA0..."
```
Mail server signs emails with private key; public key in DNS for recipient verification.

### DMARC — Domain-based Message Auth
```
_dmarc.example.com.  TXT  "v=DMARC1; p=reject; rua=mailto:dmarc@example.com"
```
`p=reject` = unauthenticated email dropped entirely + aggregate reports sent back.

---

# 04 — DNS Security

*Protecting the lookup chain — DNSSEC, DoH, DoT*

---

## DNS Security: DNSSEC, DoH, and DoT

> ⚠️ **Before:** `🔓` DNS query (UDP:53) — visible to anyone on the network
> ✅ **After:**  `🔒` Encrypted tunnel (port 443/853) — only resolver sees queries

### DNSSEC — DNS Security Extensions
- Digital signatures on all DNS records — chain of trust root → TLD → zone
- Does **NOT** encrypt traffic — just signs it to prevent tampering

### DoH — DNS over HTTPS (RFC 8484)
- DNS inside HTTPS on port 443 — looks like normal web traffic, hard to block
- Supported by Firefox, Chrome, Windows 11

### DoT — DNS over TLS (RFC 7858)
- DNS inside TLS on port 853 — visible to network admins, preferred in enterprise
- Android 9+ "Private DNS" feature uses DoT

---

## DNS Record Quick Reference

| Record | Purpose | Key Info | Example |
|--------|---------|----------|---------|
| `A` | Name → IPv4 | Most fundamental | `203.0.113.42` |
| `AAAA` | Name → IPv6 | 128-bit, dual-stack | `2001:db8::1` |
| `CNAME` | Alias → Name | No zone apex | `www → example.com.` |
| `MX` | Email routing | Lower priority = first | `10 mail.example.com.` |
| `TXT` | Arbitrary text | SPF, DKIM, DMARC | `v=spf1 include:... ~all` |
| `SOA` | Zone authority | One per zone | `ns1. admin. 20240401 ...` |
| `NS` | Authoritative servers | Min 2 recommended | `ns1.hover.com.` |
| `PTR` | IP → Name | IP block owner | `42.113.0.203.in-addr.arpa` |
| `SRV` | Service discovery | Priority, weight, port | `_xmpp._tcp 5 50 5222 host.` |
| `CAA` | CA authorization | SSL cert control | `0 issue "letsencrypt.org"` |
| `SSHFP` | SSH fingerprint | Needs DNSSEC | `2 1 7491973e5f8b39...` |
| `TLSA` | TLS cert pin | DANE validation | `3 1 1 abc123def456...` |

---

## Key Takeaways

1. **DNS is infrastructure** — Every internet connection depends on DNS.
2. **Records have specific jobs** — A/AAAA serve IPs, MX routes email, SRV discovers services.
3. **DNS powers email security** — SPF, DKIM, and DMARC are all DNS-based.
4. **DNS is increasingly secure** — DNSSEC signs records, DoH/DoT encrypt queries.
5. **TTL controls propagation** — Plan change windows around TTL values.

### Go Further
- RFC 1034/1035 — DNS specification  |  RFC 1912 — Common DNS errors
- IANA DNS Parameters  |  [MXToolbox](https://mxtoolbox.com)
- [dnsviz.net](https://dnsviz.net) — DNSSEC visualization
- Tools: `dig`, `nslookup`, `drill`  |  Pi-hole for DNS filtering

---

*Questions? — cascadesteam.org*
'''


def build_markdown():
    with open(MD_OUT, 'w') as f:
        f.write(MARKDOWN)
    print(f"✓ Markdown saved: {MD_OUT}")


if __name__ == "__main__":
    print("Building Cascade STEAM DNS Presentation — v2 (enhanced visuals)...")
    print(f"Output directory: {OUT_DIR}\n")
    build_pptx()
    build_markdown()
    print("\nDone.")
